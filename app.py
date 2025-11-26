import eventlet
eventlet.monkey_patch()

from flask import Flask, render_template, request, send_file, send_from_directory
from flask_socketio import SocketIO, emit
import os
import io
import time
import json
import re
import requests
import html
import urllib.parse
import base64
import traceback
from PIL import Image

# PPTX Imports
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# DOCX Imports
from docx import Document
from docx.shared import Inches as DocInches, Pt as DocPt, RGBColor as DocRGB
from docx.enum.text import WD_ALIGN_PARAGRAPH

# XLSX Imports
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment

# Reliable Web Search
from duckduckgo_search import DDGS

# --- CONFIG ---
app = Flask(__name__)
app.config['SECRET_KEY'] = os.environ.get('SECRET_KEY', 'secret!')
# Allow CORS for Render domains
socketio = SocketIO(app, async_mode='eventlet', cors_allowed_origins="*")

# --- GLOBAL STATE ---
active_generations = {}

# --- PROMPTS ---

# 1. PPT FAST MODE
PPT_FAST_SYSTEM_PROMPT = """You are a strictly JSON-speaking presentation engine.
Output valid JSON only. Language: {language}.
Do not output markdown code blocks. Just raw JSON.
Context: {context}.
JSON STRUCTURE:
{{
    "title": "Title", "subtitle": "Subtitle",
    "slides": [ {{"title": "Title", "bullets": ["A", "B"], "image_search_query": "query"}} ],
    "conclusion": {{"title": "Conc", "bullets": ["Sum"], "image_search_query": "query"}}
}}
"""

# 2. PPT DEEP MODE
PPT_PLANNER_PROMPT = """You are a Presentation Architect.
Your goal: Create a detailed To-Do list (Outline) for a presentation about "{topic}".
Context: {context}.
Language: {language}.
Output strictly JSON. No markdown.
Structure:
{{
    "presentation_title": "Main Title",
    "subtitle": "Subtitle",
    "todo_list": [
        "Topic for Slide 1",
        "Topic for Slide 2",
        "Topic for Slide 3",
        "Conclusion"
    ]
}}
"""

PPT_WORKER_PROMPT = """You are a Presentation Designer.
Your Task: Create ONE slide for the topic: "{slide_topic}".
Global Topic: {topic}. 
SEARCH RESULTS: {web_context}
Language: {language}.
Style: {style}.

STRICT RULES:
1. Use specific numbers, dates, and names from the SEARCH RESULTS.
2. Do not invent people or financial figures.
3. If search results are empty, write general concepts only.

Output strictly JSON. No markdown.
Structure:
{{
    "title": "Slide Title",
    "bullets": ["Real fact 1", "Real fact 2", "Real fact 3"],
    "image_search_query": "query (physical object, max 3 words)"
}}
"""

DOC_PLANNER_PROMPT = """You are a Document Architect.
Create a structural outline for a document about: "{topic}".
Context: {context}.
Language: {language}.

STRICT RULES:
1. "sections" must be a simple list of STRINGS.
2. DO NOT create objects or subsections inside the list.

CORRECT:
"sections": ["Introduction", "History", "Features"]

INCORRECT:
"sections": [{{"title": "Introduction", "subsections": [...]}}]

Output strictly JSON.
Structure:
{{
    "title": "Main Title",
    "subtitle": "Subtitle",
    "sections": [
        "Introduction to {topic}",
        "History of {topic}",
        "Key Features",
        "Modern Applications",
        "Conclusion"
    ]
}}
"""

DOC_WORKER_PROMPT = """You are a Technical Writer.
Write one detailed section for the heading: "{section_heading}".
Document Topic: {topic}.
SEARCH RESULTS: {web_context}
Language: {language}.

STRICT RULES:
1. Base your paragraph ONLY on the SEARCH RESULTS provided.
2. Output strictly valid JSON. 
3. DO NOT nest the JSON inside the heading name. 

CORRECT FORMAT:
{{
    "heading": "{section_heading}",
    "content": "Full paragraph text here...",
    "image_search_query": "visual query"
}}
"""

EXCEL_SYSTEM_PROMPT = """You are a Data Analyst. Output strictly JSON.
Language: {language}.
SEARCH RESULTS: {web_context}
Use the SEARCH RESULTS to fill the rows with real data (names, net worths, dates) where possible.
Structure:
{{
    "filename_suggestion": "Topic_Data",
    "sheets": [
        {{
            "name": "SheetName",
            "headers": ["Col1", "Col2", "Col3"],
            "rows": [
                ["Row1Val1", 100, true],
                ["Row2Val1", 200, false]
            ]
        }}
    ]
}}
"""

# --- THEMES ---
THEMES = {
    "Corporate Blue": {"bg": (255,255,255), "title": (0,32,96), "text": (50,50,50), "accent": (0,32,96)},
    "Dark Mode": {"bg": (30,30,30), "title": (255,215,0), "text": (220,220,220), "accent": (70,70,70)},
    "Forest Green": {"bg": (240,255,240), "title": (34,139,34), "text": (0,50,0), "accent": (34,139,34)},
    "Cyber Neon": {"bg": (10,10,20), "title": (0,255,255), "text": (255,0,255), "accent": (0,255,255)},
}

# --- HELPERS ---

def img_to_base64(img_bytes):
    return "data:image/png;base64," + base64.b64encode(img_bytes).decode('utf-8')

def extract_json(raw_text):
    raw_text = re.sub(r'```json', '', raw_text)
    raw_text = re.sub(r'```', '', raw_text)
    match = re.search(r'\{.*\}', raw_text, re.DOTALL)
    if not match: raise Exception("AI returned invalid JSON")
    return json.loads(match.group(0))

def call_llm(prompt):
    url = "https://apifreellm.com/api/chat"
    headers = {"Content-Type": "application/json"}
    for attempt in range(1, 4):
        try:
            res = requests.post(url, headers=headers, json={"message": prompt}, timeout=60)
            if res.status_code != 200: continue
            data = res.json()
            if "response" in data: return data["response"]
            elif "message" in data: return data["message"]
        except Exception as e:
            print(f"LLM Err: {e}")
            time.sleep(2)
    raise Exception("AI API is busy. Try again.")

def check_stop(sid):
    if not active_generations.get(sid, False):
        socketio.emit('stopped_by_user', room=sid)
        return True
    return False

def get_web_context(query):
    print(f"DEBUG: Searching text for: {query}")
    try:
        with DDGS() as ddgs:
            results = list(ddgs.text(query, max_results=5))
        if not results: return "No specific web data found."
        context_str = ""
        for r in results: context_str += f"[Title: {r['title']}] {r['body']}\n"
        return context_str
    except Exception as e:
        print(f"Search Error: {e}")
        return "Search failed."

# --- LOGIC ENGINES ---

class WebImageSearcher:
    def __init__(self, style="", search_engine="bing"):
        self.style = style
        self.search_engine = search_engine
        self.headers = {
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36"
        }

    def search(self, query, sid=None):
        if not query: return []
        q = f"{query} {self.style}".strip()
        socketio.emit('status_update', {'msg': f"Looking for images: \"{query}\"..."})
        socketio.emit('browser_type', {'text': q})

        images = []
        if self.search_engine == "bing": images = self._search_bing(q)
        else: images = self._search_duckduckgo(q)
            
        if sid and check_stop(sid): return []
        if not images: images = self._search_wikimedia(q)
        if sid and check_stop(sid): return []
        if not images:
            socketio.emit('status_update', {'msg': f"Generating AI Art for \"{query}\"..."})
            images = self._gen_pollinations(q)
        
        if images:
            b64_images = [img_to_base64(i) for i in images[:4]]
            socketio.emit('browser_results', {'images': b64_images})
            socketio.sleep(1.5)
            socketio.emit('anim_drag', {})
            socketio.sleep(1)
        return images

    def _search_bing(self, query):
        try:
            url = f"https://www.bing.com/images/search?q={urllib.parse.quote(query)}"
            res = requests.get(url, headers=self.headers, timeout=5)
            links = re.findall(r'murl&quot;:&quot;(http[^&]+?\.jpg)&quot;', res.text)
            imgs = []
            for l in links[:4]:
                c = self._down(html.unescape(l))
                if c: imgs.append(c)
            return imgs
        except: return []

    def _search_duckduckgo(self, query):
        try:
            url = f"https://duckduckgo.com/i.js?q={urllib.parse.quote(query)}"
            res = requests.get(url, headers=self.headers, timeout=8)
            if res.status_code != 200: return []
            data = res.json()
            imgs = []
            for item in data.get('results', [])[:4]:
                img_url = item.get('image')
                if img_url:
                    c = self._down(img_url)
                    if c: imgs.append(c)
            return imgs
        except: return []

    def _search_wikimedia(self, query):
        try:
            url = "https://commons.wikimedia.org/w/api.php"
            params = { "action": "query", "format": "json", "generator": "search", "gsrnamespace": 6, "gsrsearch": query, "gsrlimit": 4, "prop": "imageinfo", "iiprop": "url" }
            res = requests.get(url, params=params, headers=self.headers, timeout=5)
            data = res.json()
            pages = data.get("query", {}).get("pages", {})
            imgs = []
            for pid in pages:
                try:
                    url = pages[pid]["imageinfo"][0]["url"]
                    if url.endswith(('.jpg', '.png')):
                        c = self._down(url)
                        if c: imgs.append(c)
                except: continue
            return imgs
        except: return []

    def _gen_pollinations(self, query):
        try:
            safe = re.sub(r'[^\w\s]', '', query).replace(' ', '%20')
            url = f"https://image.pollinations.ai/prompt/{safe}?width=800&height=600&nologo=true"
            res = requests.get(url, timeout=10)
            return [res.content] if res.status_code == 200 else []
        except: return []

    def _down(self, url):
        try:
            r = requests.get(url, headers=self.headers, timeout=4)
            if r.status_code == 200 and len(r.content) > 5000: return r.content
        except: return None

# --- FILE BUILDERS ---

class WebPPTBuilder:
    def __init__(self, theme_name):
        self.prs = Presentation()
        self.theme = THEMES.get(theme_name, THEMES["Corporate Blue"])
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)

    def _set_bg(self, slide):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*self.theme["bg"])

    def add_slide(self, title, bullets, img_bytes, is_title=False):
        layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(layout)
        self._set_bg(slide)
        if img_bytes:
            try: 
                top_pos = Inches(1.5 if is_title else 2.2)
                slide.shapes.add_picture(io.BytesIO(img_bytes), Inches(7.5), top_pos, width=Inches(5.5))
            except: pass
        top = Inches(2.5) if is_title else Inches(0.6)
        tb = slide.shapes.add_textbox(Inches(0.5), top, Inches(6.5), Inches(2))
        p = tb.text_frame.paragraphs[0]
        p.text = str(title)
        p.font.size = Pt(54 if is_title else 36)
        p.font.bold = True
        p.font.color.rgb = RGBColor(*self.theme["title"])

        if bullets:
            box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5 if is_title else 2.2), Inches(6.5), Inches(5))
            tf = box.text_frame
            tf.word_wrap = True
            for b in bullets:
                p = tf.add_paragraph()
                p.text = b if is_title else f"• {b}"
                p.font.size = Pt(32 if is_title else 20)
                p.font.color.rgb = RGBColor(*self.theme["text"])

class WebDocBuilder:
    def __init__(self, theme_name):
        self.doc = Document()
        
    def add_title(self, title, subtitle):
        head = self.doc.add_heading(str(title), 0)
        head.alignment = WD_ALIGN_PARAGRAPH.CENTER
        sub = self.doc.add_paragraph(str(subtitle))
        sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
        sub.style = "Subtitle"

    def add_section(self, heading, content, img_bytes=None):
        self.doc.add_heading(str(heading), level=1) 
        if img_bytes:
            try:
                with io.BytesIO(img_bytes) as image_stream:
                    self.doc.add_picture(image_stream, width=DocInches(4))
                    last_paragraph = self.doc.paragraphs[-1] 
                    last_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
            except: pass
        p = self.doc.add_paragraph(str(content))
        p.paragraph_format.space_after = DocPt(12)

    def save(self, path):
        self.doc.save(path)

class WebSheetBuilder:
    def __init__(self):
        self.wb = Workbook()
        self.wb.remove(self.wb.active)

    def add_sheet(self, name, headers, rows):
        ws = self.wb.create_sheet(title=str(name)[:30])
        header_font = Font(bold=True, color="FFFFFF")
        header_fill = PatternFill(start_color="003366", end_color="003366", fill_type="solid")
        ws.append(headers)
        for cell in ws[1]:
            cell.font = header_font
            cell.fill = header_fill
            cell.alignment = Alignment(horizontal='center')
        for row in rows: ws.append(row)
        for col in ws.columns:
            max_length = 0
            column = col[0].column_letter
            for cell in col:
                try:
                    if len(str(cell.value)) > max_length: max_length = len(str(cell.value))
                except: pass
            ws.column_dimensions[column].width = (max_length + 2)
    def save(self, path): self.wb.save(path)

# --- FLASK ROUTES ---

@app.route('/')
def index(): return render_template('index.html', themes=THEMES.keys())

@app.route('/download/<filename>')
def download(filename): return send_file(f"generated/{filename}", as_attachment=True)

@app.route('/logo.png')
def serve_logo(): return send_from_directory('.', 'logo.png')

# --- SOCKET EVENTS ---

@socketio.on('start_generation')
def handle_generation(data):
    sid = request.sid
    active_generations[sid] = True
    socketio.start_background_task(target=process_presentation, data=data, sid=sid)

@socketio.on('stop_generation')
def handle_stop():
    sid = request.sid
    if sid in active_generations: active_generations[sid] = False

@socketio.on('disconnect')
def handle_disconnect():
    sid = request.sid
    if sid in active_generations: del active_generations[sid]

# --- MAIN LOGIC ---

def process_presentation(data, sid):
    try:
        topic = data['topic']
        lang = data['lang']
        style = data['style']
        theme = data['theme']
        deep_mode = data.get('deep_mode', False)
        file_format = data.get('format', 'pptx')
        
        if check_stop(sid): return

        socketio.emit('status_update', {'msg': f"Gathering Context for {topic}..."}, room=sid)
        global_context = get_web_context(topic)

        if check_stop(sid): return

        search_engine = data.get('search_engine', 'bing')
        searcher = WebImageSearcher(style, search_engine)

        if not os.path.exists('generated'): os.makedirs('generated')
        filename = ""

        # ==========================================
        #  PPTX MODE
        # ==========================================
        if file_format == 'pptx':
            builder = WebPPTBuilder(theme)
            
            if not deep_mode:
                # === FAST MODE ===
                socketio.emit('status_update', {'msg': "Generating Fast Draft..."}, room=sid)
                
                # 1. Generate All Content at once
                prompt = PPT_FAST_SYSTEM_PROMPT.format(language=lang, context=global_context) + f"\nRequest: {topic}"
                raw = call_llm(prompt)
                if check_stop(sid): return
                content = extract_json(raw)
                
                # 2. Add Slides
                slides = [ {'title': content['title'], 'bullets': [content.get('subtitle','')], 'q': topic+" conceptual", 'is_title': True} ]
                for s in content.get('slides', []):
                    slides.append({'title': s['title'], 'bullets': s['bullets'], 'q': s['image_search_query'], 'is_title': False})
                if 'conclusion' in content:
                    c = content['conclusion']
                    slides.append({'title': c['title'], 'bullets': c['bullets'], 'q': c['image_search_query'], 'is_title': False})

                total = len(slides)
                for i, s in enumerate(slides):
                    if check_stop(sid): return
                    socketio.emit('status_update', {'msg': f"Slide {i+1}/{total}: {s['title']}"}, room=sid)
                    socketio.emit('slide_type', {'title': s['title'], 'bullets': s['bullets']}, room=sid)

                    imgs = searcher.search(s['q'], sid)
                    if check_stop(sid): return
                    
                    builder.add_slide(s['title'], s['bullets'], imgs[0] if imgs else None, s['is_title'])
                    socketio.sleep(0.5)

            else:
                # === DEEP THINKING MODE ===
                socketio.emit('status_update', {'msg': "Designing Presentation Structure..."}, room=sid)
                plan_prompt = PPT_PLANNER_PROMPT.format(topic=topic, context=global_context, language=lang)
                raw_plan = call_llm(plan_prompt)
                if check_stop(sid): return
                plan = extract_json(raw_plan)
                
                todo_list = plan.get('todo_list', [])
                main_title = plan.get('presentation_title', topic)
                subtitle = plan.get('subtitle', '')
                
                socketio.emit('planner_init', {'items': todo_list}, room=sid)
                socketio.emit('slide_type', {'title': main_title, 'bullets': [subtitle]}, room=sid)
                
                imgs = searcher.search(f"{topic} {style}", sid)
                builder.add_slide(main_title, [subtitle], imgs[0] if imgs else None, True)
                
                total = len(todo_list)
                for i, slide_task in enumerate(todo_list):
                    if check_stop(sid): return
                    socketio.emit('status_update', {'msg': f"Fact-checking: {slide_task}..."}, room=sid)
                    specific_context = get_web_context(f"{topic} {slide_task} statistics facts")
                    
                    socketio.emit('status_update', {'msg': f"Slide {i+1}/{total}: {slide_task}"})
                    work_prompt = PPT_WORKER_PROMPT.format(slide_topic=slide_task, topic=topic, web_context=specific_context, language=lang, style=style)
                    
                    raw_slide = call_llm(work_prompt)
                    if check_stop(sid): return
                    slide_content = extract_json(raw_slide)
                    
                    s_title = slide_content.get('title', slide_task)
                    s_bullets = slide_content.get('bullets', [])
                    s_query = slide_content.get('image_search_query', slide_task)
                    
                    socketio.emit('slide_type', {'title': s_title, 'bullets': s_bullets}, room=sid)
                    imgs = searcher.search(s_query, sid)
                    if check_stop(sid): return
                    builder.add_slide(s_title, s_bullets, imgs[0] if imgs else None, False)
                    socketio.emit('mark_done', {'index': i}, room=sid)
                    socketio.sleep(0.5)
            
            filename = f"Presentation_{int(time.time())}.pptx"
            builder.prs.save(os.path.join('generated', filename))

        # ==========================================
        #  DOCX MODE
        # ==========================================
        elif file_format == 'docx':
            builder = WebDocBuilder(theme)
            socketio.emit('status_update', {'msg': f"Planning Document..."}, room=sid)
            
            prompt = DOC_PLANNER_PROMPT.format(topic=topic, context=global_context, language=lang)
            raw = call_llm(prompt)
            if check_stop(sid): return
            plan = extract_json(raw)
            
            builder.add_title(plan.get('title', topic), plan.get('subtitle', ''))
            
            # Sanitization
            raw_sections = plan.get('sections', [])
            clean_sections = []
            for s in raw_sections:
                if isinstance(s, str): clean_sections.append(s)
                elif isinstance(s, dict):
                    if 'title' in s: clean_sections.append(str(s['title']))
                    elif 'heading' in s: clean_sections.append(str(s['heading']))
                    elif 'name' in s: clean_sections.append(str(s['name']))
                    else: 
                        for v in s.values():
                            if isinstance(v, str): clean_sections.append(v); break
                else: clean_sections.append(str(s))

            socketio.emit('planner_init', {'items': clean_sections}, room=sid)
            total = len(clean_sections)

            for i, sec_name in enumerate(clean_sections):
                if check_stop(sid): return
                
                socketio.emit('status_update', {'msg': f"Researching: {sec_name}..."}, room=sid)
                local_context = get_web_context(f"{topic} {sec_name} detailed facts")

                socketio.emit('status_update', {'msg': f"Writing Section {i+1}/{total}..."})
                w_prompt = DOC_WORKER_PROMPT.format(section_heading=sec_name, topic=topic, web_context=local_context, language=lang)
                
                raw_sec = call_llm(w_prompt)
                content = extract_json(raw_sec)

                final_heading = content.get('heading', sec_name)
                final_content = content.get('content', '')
                if final_content == "" and final_heading in content and isinstance(content[final_heading], dict):
                     sub = content[final_heading]
                     final_content = sub.get('content', '')

                preview_text = final_content[:150] + "..."
                socketio.emit('slide_type', {'title': final_heading, 'bullets': [preview_text]}, room=sid)
                
                imgs = searcher.search(content.get('image_search_query', topic), sid)
                
                builder.add_section(final_heading, final_content, imgs[0] if imgs else None)
                socketio.emit('mark_done', {'index': i}, room=sid)
                socketio.sleep(0.5)
            
            filename = f"Document_{int(time.time())}.docx"
            builder.save(os.path.join('generated', filename))

        # ==========================================
        #  XLSX MODE
        # ==========================================
        elif file_format == 'xlsx':
            builder = WebSheetBuilder()
            socketio.emit('status_update', {'msg': f"Searching Data Sources..."}, room=sid)
            
            data_context = get_web_context(f"{topic} statistics data table net worth list")
            if len(data_context) < 50: data_context = get_web_context(f"top 10 {topic} list stats")

            prompt = EXCEL_SYSTEM_PROMPT.format(language=lang, web_context=data_context) + f"\nRequest: {topic}"
            
            raw = call_llm(prompt)
            if check_stop(sid): return
            content = extract_json(raw)
            
            sheets = content.get('sheets', [])
            total = len(sheets)
            
            sheet_names = [s['name'] for s in sheets]
            socketio.emit('planner_init', {'items': sheet_names}, room=sid)

            for i, sheet in enumerate(sheets):
                if check_stop(sid): return
                socketio.emit('status_update', {'msg': f"Building Sheet {i+1}/{total}: {sheet['name']}"}, room=sid)
                socketio.emit('slide_type', {'title': sheet['name'], 'bullets': sheet['headers']}, room=sid)
                
                builder.add_sheet(sheet['name'], sheet['headers'], sheet['rows'])
                socketio.emit('mark_done', {'index': i}, room=sid)
                socketio.sleep(1)
            
            filename = f"Spreadsheet_{int(time.time())}.xlsx"
            builder.save(os.path.join('generated', filename))

        # === FINISH ===
        socketio.emit('finished', {'filename': filename}, room=sid)
        if sid in active_generations: del active_generations[sid]

    except Exception as e:
        print(f"CRITICAL ERROR: {e}")
        traceback.print_exc()
        socketio.emit('error', {'msg': str(e)}, room=sid)
        if sid in active_generations: del active_generations[sid]

if __name__ == '__main__':
    socketio.run(app, host='0.0.0.0', port=5000)
